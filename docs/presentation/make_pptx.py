# Build presentation.pptx as a native-shape replica of the beamer PDF.
# Geometry, colors and fonts are extracted from the PDF itself; speaker
# notes come from the \note{} blocks in presentation.tex.
import re
import sys
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
PDF = str(HERE / "build" / "presentation.pdf")
TEX = str(HERE / "presentation.tex")
OUT = sys.argv[1] if len(sys.argv) > 1 else str(HERE / "build" / "presentation.pptx")

SLIDE_W_EMU = 9144000          # 10 in
SLIDE_H_EMU = 6858000          # 7.5 in

doc = fitz.open(PDF)
PAGE_W = doc[0].rect.width     # 362.83 pt
PAGE_H = doc[0].rect.height    # 272.13 pt
S = SLIDE_W_EMU / PAGE_W       # EMU per PDF pt
FS = S / 12700.0               # font size multiplier (pdf pt -> ppt pt)

# Calibration knobs (tuned via render-compare loop)
BASELINE_K = 1.03              # first-baseline offset as fraction of font size
BASELINE_K_MATH = 0.81         # LM Math sets USE_TYPO_METRICS (typoAsc 806/1000)
X_NUDGE_EMU = 0                # horizontal fudge

# ---------------------------------------------------------------- notes ----
def extract_notes(tex_path):
    src = open(tex_path, encoding="utf-8").read()
    notes = []
    i = 0
    while True:
        m = re.search(r"\\note\{", src[i:])
        if not m:
            break
        start = i + m.end()
        depth = 1
        j = start
        while depth:
            if src[j] == "{":
                depth += 1
            elif src[j] == "}":
                depth -= 1
            j += 1
        notes.append(src[start:j - 1])
        i = j
    def clean(t):
        t = t.strip()
        t = t.replace("---", "—")
        t = t.replace(",,", "„").replace("''", "”")
        t = t.replace("$<$", "<").replace(r"\&", "&").replace(r"\%", "%")
        paras = re.split(r"\n\s*\n", t)
        paras = [re.sub(r"\s*\n\s*", " ", p).strip() for p in paras]
        return "\n\n".join(paras)
    return [clean(n) for n in notes]

NOTES = extract_notes(TEX)

# ------------------------------------------------------------- helpers ----
def rgb(col):
    return RGBColor(*(round(c * 255) for c in col))

def set_no_style(shape):
    """Remove theme style reference so explicit fill/line rule."""
    el = shape._element.find(qn("p:style"))
    if el is not None:
        shape._element.remove(el)

def add_shadow(shape, blur_pt=4.0, dist_pt=2.6, dir_deg=45, alpha=35):
    spPr = shape._element.spPr
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shd = spPr.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur_pt * 12700 * (S / 12700.0))),
        "dist": str(int(dist_pt * S)),
        "dir": str(int(dir_deg * 60000)),
        "rotWithShape": "0",
    })
    clr = spPr.makeelement(qn("a:srgbClr"), {"val": "000000"})
    a = spPr.makeelement(qn("a:alpha"), {"val": str(alpha * 1000)})
    clr.append(a)
    shd.append(clr)
    eff.append(shd)
    spPr.append(eff)

def add_tail_arrow(shape, kind="triangle"):
    ln = shape.line._get_or_add_ln()
    end = ln.makeelement(qn("a:tailEnd"), {"type": kind, "w": "lg", "len": "lg"})
    ln.append(end)

def corner_radius(d):
    """Estimate rounded-corner radius from curve endpoints vs bbox corners."""
    r = d["rect"]
    corners = [(r.x0, r.y0), (r.x1, r.y0), (r.x0, r.y1), (r.x1, r.y1)]
    dists = []
    for it in d["items"]:
        if it[0] != "c":
            continue
        for p in (it[1], it[4]):
            dmin = min(max(abs(p.x - cx), abs(p.y - cy)) for cx, cy in corners)
            dists.append(dmin)
    if not dists:
        return 0.0
    dists.sort()
    return dists[len(dists) // 2]

def font_map(pdf_font, text=""):
    """Map PDF font (with optical size) to an installed family."""
    bold = "Bold" in pdf_font
    italic = "Oblique" in pdf_font or "Italic" in pdf_font
    if pdf_font.startswith(("CM", "LatinModernMath")) and text.strip() not in ("<", ">"):
        return "Latin Modern Math", False, False
    if "LMSans12" in pdf_font:
        return "LM Sans 12", bold, italic
    if "LMSans8" in pdf_font:
        return "LM Sans 8", bold, italic
    if "LMSans9" in pdf_font:
        # glyphs already oblique in this face; don't add faux italic
        return "LM Sans 9 Oblique", bold, False
    return "Latin Modern Sans", bold, italic

# ------------------------------------------------------------- builder ----
prs = Presentation()
prs.slide_width = Emu(SLIDE_W_EMU)
prs.slide_height = Emu(SLIDE_H_EMU)
blank = prs.slide_layouts[6]

def add_rect(shapes, r, fill, stroke=None, width=None, radius_pt=0.0, shadow=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_pt > 0.2 else MSO_SHAPE.RECTANGLE
    shp = shapes.add_shape(shp_type,
                           Emu(int(r.x0 * S)), Emu(int(r.y0 * S)),
                           Emu(int((r.x1 - r.x0) * S)), Emu(int((r.y1 - r.y0) * S)))
    set_no_style(shp)
    if radius_pt > 0.2:
        shp.adjustments[0] = radius_pt / min(r.x1 - r.x0, r.y1 - r.y0)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
    else:
        shp.fill.background()
    if stroke:
        shp.line.color.rgb = rgb(stroke)
        shp.line.width = Emu(int(width * S))
    else:
        shp.line.fill.background()
    if shadow:
        add_shadow(shp)
    shp.shadow.inherit = False
    return shp

def add_line(shapes, p1, p2, color, width, arrow=False):
    conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Emu(int(p1.x * S)), Emu(int(p1.y * S)),
                                Emu(int(p2.x * S)), Emu(int(p2.y * S)))
    set_no_style(conn)
    conn.line.color.rgb = rgb(color)
    conn.line.width = Emu(int(width * S))
    conn.shadow.inherit = False
    if arrow:
        add_tail_arrow(conn)
    return conn

def span_words(span):
    """Split a rawdict span into words with exact per-word x/baseline.

    Words break on space chars AND on gaps wider than ~40% of the font
    size (TeX inter-word glue has no space glyph)."""
    size = span["size"]
    words = []
    cur = ""
    cur_x = cur_y = None
    prev_x1 = None
    for ch in span["chars"]:
        c = ch["c"]
        x0 = ch["origin"][0]
        if c == " ":
            if cur:
                words.append((cur, cur_x, cur_y))
                cur = ""
            prev_x1 = None
            continue
        if cur and prev_x1 is not None and x0 - prev_x1 > 0.4 * size:
            words.append((cur, cur_x, cur_y))
            cur = ""
        if not cur:
            cur_x, cur_y = ch["origin"]
        cur += c
        prev_x1 = ch["bbox"][2]
    if cur:
        words.append((cur, cur_x, cur_y))
    return words

def add_text(shapes, span):
    size_pt = span["size"] * FS
    boxes = []
    for text, wx, wy in span_words(span):
        name, bold, italic = font_map(span["font"], text)
        k = BASELINE_K_MATH if name == "Latin Modern Math" else BASELINE_K
        x = wx * S + X_NUDGE_EMU
        top = wy * S - k * size_pt * 12700
        w = int(len(text) * size_pt * 12700) + int(8 * S)
        h = int(size_pt * 1.5 * 12700)
        tb = shapes.add_textbox(Emu(int(x)), Emu(int(top)), Emu(w), Emu(h))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = name
        f.size = Pt(size_pt)
        f.bold = bold
        f.italic = italic
        col = span["color"]
        f.color.rgb = RGBColor((col >> 16) & 255, (col >> 8) & 255, col & 255)
        boxes.append(tb)
    return boxes

def find_ball_bullets(page):
    """Beamer ball bullets are shadings (invisible to get_drawings).
    Locate them by scanning rendered pixels in the left text margin,
    away from the frametitle bar and the footline."""
    pix = page.get_pixmap(dpi=150)
    sc = 150.0 / 72.0
    x_lo, x_hi = int(15 * sc), int(31 * sc)
    y_lo, y_hi = int(35 * sc), int(255 * sc)

    def is_ball_px(p):
        return p[0] + p[1] + p[2] < 690   # any non-white pixel in bullet gutter

    balls = []
    y = y_lo
    while y < y_hi:
        hit0 = [x for x in range(x_lo, x_hi) if is_ball_px(pix.pixel(x, y))]
        if hit0:
            y0 = y
            xs = []
            while y < y_hi:
                hit = [x for x in range(x_lo, x_hi) if is_ball_px(pix.pixel(x, y))]
                if not hit:
                    break
                xs.extend(hit)
                y += 1
            h_pt = (y - y0) / sc
            w_pt = (max(xs) - min(xs) + 1) / sc
            if 2.0 <= h_pt <= 7.0 and w_pt <= 7.0:   # bullet-sized blob only
                cx = (min(xs) + max(xs) + 1) / 2 / sc
                cy = (y0 + y) / 2 / sc
                balls.append((cx, cy, w_pt / 2))
        y += 1
    return balls

def add_ball(shapes, cx, cy, r):
    shp = shapes.add_shape(MSO_SHAPE.OVAL,
                           Emu(int((cx - r) * S)), Emu(int((cy - r) * S)),
                           Emu(int(2 * r * S)), Emu(int(2 * r * S)))
    set_no_style(shp)
    shp.line.fill.background()
    shp.shadow.inherit = False
    # radial gradient: light center -> structure blue rim (beamer ball look)
    spPr = shp._element.spPr
    fill_el = spPr.find(qn("a:solidFill"))
    if fill_el is not None:
        spPr.remove(fill_el)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    lst = spPr.makeelement(qn("a:gsLst"), {})
    for pos, val in ((0, "9999E6"), (35000, "4D4DB8"), (100000, "26268C")):
        gs = spPr.makeelement(qn("a:gs"), {"pos": str(pos)})
        c = spPr.makeelement(qn("a:srgbClr"), {"val": val})
        gs.append(c)
        lst.append(gs)
    grad.append(lst)
    path = spPr.makeelement(qn("a:path"), {"path": "circle"})
    rect = spPr.makeelement(qn("a:fillToRect"), {"l": "30000", "t": "30000",
                                                 "r": "70000", "b": "70000"})
    path.append(rect)
    grad.append(path)
    ln = spPr.find(qn("a:ln"))
    spPr.insert(list(spPr).index(ln) if ln is not None else len(list(spPr)), grad)
    return shp

for pno, page in enumerate(doc):
    slide = prs.slides.add_slide(blank)
    shapes = slide.shapes

    drawings = page.get_drawings()
    # collect stroke-lines and arrowheads first so arrowheads can extend lines
    lines = []      # dicts: p1, p2, color, width, arrow(False/'p1'/'p2')
    arrowheads = [] # (center, tip candidates bbox, color)
    boxes = []      # rect-ish fills in painter order

    for d in drawings:
        r = d["rect"]
        # page background
        if d["type"] == "f" and d["fill"] == (1.0, 1.0, 1.0) and r.width > PAGE_W - 2:
            continue
        # beamer title-block fake shadows (black rects behind the block)
        if pno == 0 and d["fill"] == (0.0, 0.0, 0.0):
            continue
        if d["type"] == "s" and len(d["items"]) == 1 and d["items"][0][0] == "l":
            lines.append({"p1": d["items"][0][1], "p2": d["items"][0][2],
                          "color": d["color"], "width": d["width"], "arrow": None})
            continue
        if d["type"] == "f" and len(d["items"]) == 3 and r.width < 12 and r.height < 12:
            arrowheads.append(d)
            continue
        boxes.append(d)

    # page 1: merge the two blue pieces of the rounded title block
    if pno == 0:
        blues = [d for d in boxes if d["fill"] and abs(d["fill"][2] - 0.7) < 0.01
                 and d["fill"][0] < 0.25 and d["rect"].width > 300]
        if len(blues) == 2:
            r_all = fitz.Rect(min(b["rect"].x0 for b in blues),
                              min(b["rect"].y0 for b in blues),
                              max(b["rect"].x1 for b in blues),
                              max(b["rect"].y1 for b in blues))
            keep = blues[0]
            keep["rect"] = r_all
            keep["_shadow"] = True
            keep["_radius"] = 4.0
            boxes.remove(blues[1])

    # arrowheads become explicit rotated triangles (PowerPoint distorts
    # tailEnd arrows on tikz's shortened stub lines)
    triangles = []
    for ah in arrowheads:
        c = ((ah["rect"].x0 + ah["rect"].x1) / 2, (ah["rect"].y0 + ah["rect"].y1) / 2)
        best, which, bd = None, None, 1e9
        for ln_ in lines:
            for key in ("p1", "p2"):
                p = ln_[key]
                dd = (p.x - c[0]) ** 2 + (p.y - c[1]) ** 2
                if dd < bd:
                    bd, best, which = dd, ln_, key
        if best is None or bd > 150:
            continue
        # direction of travel = from other endpoint toward this one
        other = best["p2"] if which == "p1" else best["p1"]
        p = best[which]
        vx, vy = p.x - other.x, p.y - other.y
        n = (vx ** 2 + vy ** 2) ** 0.5 or 1.0
        vx, vy = vx / n, vy / n
        import math
        rot = math.degrees(math.atan2(vy, vx)) + 90.0
        L = max(ah["rect"].width, ah["rect"].height)
        W = 0.85 * L
        triangles.append((c[0], c[1], W, L, rot, ah["fill"]))

    # draw boxes in painter order
    for d in boxes:
        radius = d.get("_radius", corner_radius(d))
        fill = d["fill"] if d["type"] in ("f", "fs") else None
        stroke = d["color"] if d["type"] in ("s", "fs") else None
        add_rect(shapes, d["rect"], fill, stroke, d.get("width"),
                 radius_pt=radius, shadow=d.get("_shadow", False))

    # lines (under text, over boxes — matches PDF painter order closely enough)
    for ln_ in lines:
        add_line(shapes, ln_["p1"], ln_["p2"], ln_["color"], ln_["width"])

    # arrowhead triangles
    for cx, cy, W, L, rot, fill in triangles:
        tri = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                               Emu(int((cx - W / 2) * S)), Emu(int((cy - L / 2) * S)),
                               Emu(int(W * S)), Emu(int(L * S)))
        set_no_style(tri)
        tri.rotation = rot
        tri.fill.solid()
        tri.fill.fore_color.rgb = rgb(fill)
        tri.line.fill.background()
        tri.shadow.inherit = False

    # ball bullets (page 9 itemize)
    if pno == 8:
        for cx, cy, rr in find_ball_bullets(page):
            add_ball(shapes, cx, cy, rr)

    # text spans (rawdict: per-char coords -> exact per-word placement)
    for block in page.get_text("rawdict")["blocks"]:
        if block["type"]:
            continue
        for line in block["lines"]:
            for span in line["spans"]:
                if not span["chars"]:
                    continue
                add_text(shapes, span)

    # speaker notes
    if pno < len(NOTES):
        slide.notes_slide.notes_text_frame.text = NOTES[pno]

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
